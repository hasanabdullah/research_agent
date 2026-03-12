"""Slide/presentation generation from research files.

Reads research markdown, uses LLM to generate structured slide content,
renders to .pptx via python-pptx. Output importable into Google Slides.
"""

import json
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from agent import ROOT, load_config, load_identity
from llm import get_client, resolve_model, completions_with_retry


# --- Slide content generation via LLM ---

SLIDE_SYSTEM_PROMPT = """\
You are a presentation designer. Given research markdown files, produce a structured JSON array of slides.

Each slide is an object with these fields:
- "layout": one of "title", "content", "table", "two_column"
- "title": slide title (string)
- "bullets": list of bullet strings (for "content" layout; omit for others)
- "left_bullets": list of strings (for "two_column" layout only)
- "right_bullets": list of strings (for "two_column" layout only)
- "left_title": string (for "two_column" layout only)
- "right_title": string (for "two_column" layout only)
- "table": { "headers": [...], "rows": [[...], ...] } (for "table" layout only)
- "subtitle": string (for "title" layout only)
- "speaker_notes": string (optional, for any layout)

Rules:
- First slide MUST be layout "title" with the topic name and a subtitle
- Last slide should be recommendations / next steps
- Use BUILD or KILL verdicts prominently where applicable
- Include real numbers: competitor counts, pricing, TAM estimates, scores
- Keep bullet text concise (max ~15 words per bullet)
- Tables should have max 5-6 columns and max 8 rows
- For "concise" style: ~18-22 slides, investor-pitch feel, big numbers
- For "detailed" style: ~30-40 slides, full reference deck with per-idea breakdowns
- Return ONLY valid JSON array, no markdown fences or extra text
"""


def generate_slide_content(topic_name: str, style: str = "concise") -> list[dict]:
    """Read research files, send to LLM, return structured slide content."""
    base = ROOT / "topics" / topic_name
    research_dir = base / "data" / "research"
    if not research_dir.exists():
        raise FileNotFoundError(f"No research directory for topic '{topic_name}'")

    # Gather research files
    files = sorted(
        list(research_dir.glob("*.md")) + list(research_dir.glob("*.txt")),
        key=lambda f: f.name,
    )
    parts = []
    for f in files:
        content = f.read_text(encoding="utf-8").strip()
        if content:
            parts.append(f"## {f.name}\n\n{content}")

    if not parts:
        raise ValueError("No research content found")

    combined = "\n\n---\n\n".join(parts)

    # Load identity for topic name/purpose
    identity = {}
    id_path = base / "identity.json"
    if id_path.exists():
        identity = json.loads(id_path.read_text(encoding="utf-8"))

    topic_label = identity.get("purpose", topic_name)

    user_msg = (
        f"Topic: {topic_name}\n"
        f"Purpose: {topic_label}\n"
        f"Style: {style}\n"
        f"Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}\n\n"
        f"Research content:\n\n{combined}"
    )

    config = load_config()
    model = config.get("scaffold_model", "anthropic/claude-sonnet-4.6")
    client = get_client()

    response = completions_with_retry(
        client,
        model=resolve_model(model),
        max_tokens=8192,
        messages=[
            {"role": "system", "content": SLIDE_SYSTEM_PROMPT},
            {"role": "user", "content": user_msg},
        ],
    )

    raw = response.choices[0].message.content.strip()
    # Strip markdown fences if LLM wraps them
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1] if "\n" in raw else raw[3:]
        if raw.endswith("```"):
            raw = raw[:-3]
        raw = raw.strip()

    slides = json.loads(raw)

    # Record cost
    usage = response.usage
    if usage:
        from costs import record_call
        pricing = config.get("pricing", {"input_per_mtok": 1.0, "output_per_mtok": 5.0})
        costs_file = base / "data" / "costs.json"
        record_call(
            getattr(usage, "prompt_tokens", 0) or 0,
            getattr(usage, "completion_tokens", 0) or 0,
            pricing, label="slide_generation", costs_file=costs_file,
        )

    return slides


# --- PPTX rendering ---

# Color palette
_BG = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
_ACCENT = RGBColor(0x00, 0xD4, 0xAA)  # teal accent
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY = RGBColor(0xAA, 0xAA, 0xBB)
_DARK_BG = RGBColor(0x16, 0x16, 0x28)
_BUILD_GREEN = RGBColor(0x00, 0xCC, 0x66)
_KILL_RED = RGBColor(0xFF, 0x44, 0x44)


def _set_slide_bg(slide, color=_BG):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=_WHITE, alignment=PP_ALIGN.LEFT):
    """Add a styled textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_bullets(slide, left, top, width, height, bullets, font_size=13, color=_WHITE):
    """Add a bulleted list textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Highlight BUILD/KILL verdicts
        if "BUILD" in bullet.upper():
            p.font.color.rgb = _BUILD_GREEN
        elif "KILL" in bullet.upper():
            p.font.color.rgb = _KILL_RED
        else:
            p.font.color.rgb = color
        p.text = f"  {bullet}"
        p.font.size = Pt(font_size)
        p.space_after = Pt(4)
    return txBox


def _render_title_slide(prs, slide_data):
    """Render a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_bg(slide)

    # Title
    _add_textbox(
        slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5),
        slide_data.get("title", ""), font_size=36, bold=True,
        color=_WHITE, alignment=PP_ALIGN.CENTER,
    )
    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, Inches(1.0), Inches(3.5), Inches(8.0), Inches(1.0),
            subtitle, font_size=18, color=_ACCENT, alignment=PP_ALIGN.CENTER,
        )
    # Date line
    _add_textbox(
        slide, Inches(1.0), Inches(4.5), Inches(8.0), Inches(0.5),
        datetime.now(timezone.utc).strftime("%B %d, %Y"),
        font_size=12, color=_GRAY, alignment=PP_ALIGN.CENTER,
    )
    # Speaker notes
    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _render_content_slide(prs, slide_data):
    """Render a content slide with title + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title bar
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7),
        slide_data.get("title", ""), font_size=24, bold=True, color=_ACCENT,
    )
    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.0), Inches(2.0), Pt(3),  # 1 = rectangle
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _ACCENT
    line.line.fill.background()

    # Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        _add_bullets(
            slide, Inches(0.6), Inches(1.2), Inches(8.5), Inches(5.5),
            bullets, font_size=15,
        )

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _render_table_slide(prs, slide_data):
    """Render a table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6),
        slide_data.get("title", ""), font_size=22, bold=True, color=_ACCENT,
    )

    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        return

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    # Calculate table dimensions
    table_width = Inches(9.0)
    table_height = Inches(min(0.4 * n_rows, 5.5))
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.1), table_width, table_height,
    ).table

    # Style header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = str(h)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = _BG
        cell.fill.solid()
        cell.fill.fore_color.rgb = _ACCENT

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = _WHITE
                # Color BUILD/KILL
                if "BUILD" in str(val).upper():
                    paragraph.font.color.rgb = _BUILD_GREEN
                elif "KILL" in str(val).upper():
                    paragraph.font.color.rgb = _KILL_RED
            cell.fill.solid()
            cell.fill.fore_color.rgb = _DARK_BG if i % 2 == 0 else _BG

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


def _render_two_column_slide(prs, slide_data):
    """Render a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Title
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6),
        slide_data.get("title", ""), font_size=22, bold=True, color=_ACCENT,
    )

    # Left column header
    left_title = slide_data.get("left_title", "")
    if left_title:
        _add_textbox(
            slide, Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.5),
            left_title, font_size=16, bold=True, color=_WHITE,
        )

    # Right column header
    right_title = slide_data.get("right_title", "")
    if right_title:
        _add_textbox(
            slide, Inches(5.3), Inches(1.1), Inches(4.2), Inches(0.5),
            right_title, font_size=16, bold=True, color=_WHITE,
        )

    # Left bullets
    left_bullets = slide_data.get("left_bullets", [])
    if left_bullets:
        _add_bullets(
            slide, Inches(0.5), Inches(1.6), Inches(4.2), Inches(5.0),
            left_bullets, font_size=13,
        )

    # Right bullets
    right_bullets = slide_data.get("right_bullets", [])
    if right_bullets:
        _add_bullets(
            slide, Inches(5.3), Inches(1.6), Inches(4.2), Inches(5.0),
            right_bullets, font_size=13,
        )

    # Divider line
    divider = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.1), Pt(2), Inches(5.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
    divider.line.fill.background()

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]


_RENDERERS = {
    "title": _render_title_slide,
    "content": _render_content_slide,
    "table": _render_table_slide,
    "two_column": _render_two_column_slide,
}


def render_pptx(slides: list[dict], output_path: str) -> str:
    """Render slide content JSON to a .pptx file."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        layout = slide_data.get("layout", "content")
        renderer = _RENDERERS.get(layout, _render_content_slide)
        renderer(prs, slide_data)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# --- Orchestrator ---

def generate_presentation(topic_name: str, style: str = "concise") -> dict:
    """Full pipeline: generate slide content via LLM, render to .pptx, save both.

    Returns dict with file_path, slides_path, slide_count, style.
    """
    slides = generate_slide_content(topic_name, style=style)

    data_dir = ROOT / "topics" / topic_name / "data"
    data_dir.mkdir(parents=True, exist_ok=True)

    # Save slide content JSON (reusable for Notion publish)
    slides_json_path = data_dir / "slides.json"
    slides_json_path.write_text(
        json.dumps({
            "slides": slides,
            "style": style,
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "slide_count": len(slides),
        }, indent=2),
        encoding="utf-8",
    )

    # Render .pptx
    pptx_path = data_dir / "presentation.pptx"
    render_pptx(slides, str(pptx_path))

    return {
        "file_path": str(pptx_path),
        "slides_json_path": str(slides_json_path),
        "slide_count": len(slides),
        "style": style,
        "generated_at": datetime.now(timezone.utc).isoformat(),
    }


def slides_to_notion_blocks(slides: list[dict]) -> list[dict]:
    """Convert slide content JSON to Notion blocks for publishing.

    Each slide becomes: H2 heading + bullets/table + divider.
    Speaker notes rendered as callout blocks.
    """
    blocks = []

    for slide_data in slides:
        layout = slide_data.get("layout", "content")
        title = slide_data.get("title", "Untitled Slide")

        # Slide title as H2
        blocks.append({
            "object": "block",
            "type": "heading_2",
            "heading_2": {"rich_text": [{"type": "text", "text": {"content": title}}]},
        })

        if layout == "title":
            # Subtitle as paragraph
            subtitle = slide_data.get("subtitle", "")
            if subtitle:
                blocks.append({
                    "object": "block",
                    "type": "paragraph",
                    "paragraph": {"rich_text": [{"type": "text", "text": {"content": subtitle}}]},
                })

        elif layout == "content":
            for bullet in slide_data.get("bullets", []):
                blocks.append({
                    "object": "block",
                    "type": "bulleted_list_item",
                    "bulleted_list_item": {"rich_text": [{"type": "text", "text": {"content": bullet}}]},
                })

        elif layout == "table":
            table_data = slide_data.get("table", {})
            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            if headers:
                n_cols = len(headers)
                all_rows = [headers] + rows
                table_block = {
                    "object": "block",
                    "type": "table",
                    "table": {
                        "table_width": n_cols,
                        "has_column_header": True,
                        "has_row_header": False,
                        "children": [],
                    },
                }
                for row in all_rows:
                    cells = []
                    for j in range(n_cols):
                        val = str(row[j]) if j < len(row) else ""
                        cells.append([{"type": "text", "text": {"content": val}}])
                    table_block["table"]["children"].append({
                        "object": "block",
                        "type": "table_row",
                        "table_row": {"cells": cells},
                    })
                blocks.append(table_block)

        elif layout == "two_column":
            left_title = slide_data.get("left_title", "")
            right_title = slide_data.get("right_title", "")
            if left_title:
                blocks.append({
                    "object": "block",
                    "type": "heading_3",
                    "heading_3": {"rich_text": [{"type": "text", "text": {"content": left_title}}]},
                })
            for bullet in slide_data.get("left_bullets", []):
                blocks.append({
                    "object": "block",
                    "type": "bulleted_list_item",
                    "bulleted_list_item": {"rich_text": [{"type": "text", "text": {"content": bullet}}]},
                })
            if right_title:
                blocks.append({
                    "object": "block",
                    "type": "heading_3",
                    "heading_3": {"rich_text": [{"type": "text", "text": {"content": right_title}}]},
                })
            for bullet in slide_data.get("right_bullets", []):
                blocks.append({
                    "object": "block",
                    "type": "bulleted_list_item",
                    "bulleted_list_item": {"rich_text": [{"type": "text", "text": {"content": bullet}}]},
                })

        # Speaker notes as callout
        notes = slide_data.get("speaker_notes", "")
        if notes:
            blocks.append({
                "object": "block",
                "type": "callout",
                "callout": {
                    "rich_text": [{"type": "text", "text": {"content": notes}}],
                    "icon": {"type": "emoji", "emoji": "\U0001f4ac"},
                    "color": "gray_background",
                },
            })

        # Divider between slides
        blocks.append({"object": "block", "type": "divider", "divider": {}})

    return blocks
